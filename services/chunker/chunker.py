"""Core chunking, embedding, and upsert logic for the RAG ingestion pipeline."""

from __future__ import annotations

import hashlib
import tempfile
import time
from typing import Any

import structlog
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from ebooklib import ITEM_DOCUMENT, epub
from google.cloud import aiplatform, storage
from google.cloud.aiplatform.matching_engine import MatchingEngineIndex
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from pypdf import PdfReader
from vertexai.language_models import TextEmbeddingInput, TextEmbeddingModel

logger = structlog.get_logger()

# text-embedding-004 uses tiktoken cl100k_base tokenizer
CHUNK_SIZE = 500       # tokens
CHUNK_OVERLAP = 50     # tokens
EMBEDDING_MODEL = "text-embedding-004"
EMBEDDING_DIMENSIONS = 768
EMBEDDING_BATCH_SIZE = 20  # ~500 tokens/chunk × 20 = ~10k tokens, within 20k limit


def process_document(
    *,
    project_id: str,
    region: str,
    bucket_name: str,
    object_name: str,
    index_id: str,
    index_endpoint_id: str,
    content_type: str = "",
) -> dict[str, Any]:
    """Download a document from GCS, chunk it, embed chunks, and upsert to Vector Search.

    Args:
        project_id: GCP project ID.
        region: GCP region.
        bucket_name: Source GCS bucket name.
        object_name: GCS object key.
        index_id: Vertex AI Vector Search index ID.
        index_endpoint_id: Vertex AI Vector Search index endpoint ID.
        content_type: MIME type of the document.

    Returns:
        Dict with chunk_count, embed_latency_ms, upsert_latency_ms, total_latency_ms.
    """
    t_start = time.monotonic()

    aiplatform.init(project=project_id, location=region)

    # ── Download from GCS ───────────────────────────────────────────────
    local_path = _download_to_tempfile(bucket_name, object_name)

    # ── Chunk (branch by content type) ──────────────────────────────────
    if content_type == "application/epub+zip":
        chunk_records = _process_epub(local_path, object_name)
    elif content_type == "application/pdf":
        chunk_records = _process_pdf(local_path, object_name)
    elif content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        chunk_records = _process_docx(local_path, object_name)
    elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        chunk_records = _process_pptx(local_path, object_name)
    else:
        text = _read_text(local_path)
        chunks = _split_text(text)
        chunk_records = [
            {"text": chunk, "raw_id": f"{object_name}::chunk_{i}", "restricts": [
                {"namespace": "source", "allow_list": [object_name]},
            ]}
            for i, chunk in enumerate(chunks)
        ]

    if not chunk_records:
        logger.info("no_chunks_produced", object_name=object_name)
        elapsed = int((time.monotonic() - t_start) * 1000)
        return {
            "chunk_count": 0,
            "embed_latency_ms": 0,
            "upsert_latency_ms": 0,
            "total_latency_ms": elapsed,
        }

    texts = [r["text"] for r in chunk_records]

    # ── Embed ───────────────────────────────────────────────────────────
    t_embed = time.monotonic()
    embeddings = _embed_chunks(texts)
    embed_latency_ms = int((time.monotonic() - t_embed) * 1000)

    # ── Upsert to Vector Search ─────────────────────────────────────────
    t_upsert = time.monotonic()
    datapoints = _build_datapoints(chunk_records, embeddings)
    _upsert_vectors(index_id, datapoints)
    upsert_latency_ms = int((time.monotonic() - t_upsert) * 1000)

    total_latency_ms = int((time.monotonic() - t_start) * 1000)

    return {
        "chunk_count": len(chunk_records),
        "embed_latency_ms": embed_latency_ms,
        "upsert_latency_ms": upsert_latency_ms,
        "total_latency_ms": total_latency_ms,
    }


def _download_to_tempfile(bucket_name: str, object_name: str) -> str:
    """Download a GCS object to a temporary file and return its path."""
    client = storage.Client()
    bucket = client.bucket(bucket_name)
    blob = bucket.blob(object_name)

    suffix = "." + object_name.rsplit(".", 1)[-1] if "." in object_name else ".tmp"
    tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
    blob.download_to_filename(tmp.name)
    return tmp.name


def _read_text(path: str) -> str:
    """Read a file as UTF-8 text."""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _process_epub(local_path: str, object_name: str) -> list[dict[str, Any]]:
    """Extract chapters from an EPUB and chunk each independently.

    Args:
        local_path: Path to the downloaded EPUB file.
        object_name: GCS object key (used for datapoint IDs and metadata).

    Returns:
        List of chunk records with text, raw_id, and restricts.
    """
    book = epub.read_epub(local_path, options={"ignore_ncx": True})
    book_title = book.get_metadata("DC", "title")
    book_title = book_title[0][0] if book_title else object_name

    splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
        encoding_name="cl100k_base",
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
    )

    chunk_records: list[dict[str, Any]] = []
    chapter_index = 0

    for item in book.get_items_of_type(ITEM_DOCUMENT):
        html_content = item.get_content().decode("utf-8", errors="replace")
        soup = BeautifulSoup(html_content, "html.parser")

        text = soup.get_text(separator="\n", strip=True)
        if not text.strip():
            continue

        # Extract chapter title from the first heading, fall back to filename
        heading = soup.find(["h1", "h2", "h3"])
        chapter_title = heading.get_text(strip=True) if heading else item.get_name()

        chunks = splitter.split_text(text)

        for i, chunk in enumerate(chunks):
            raw_id = f"{object_name}::ch{chapter_index}::chunk_{i}"
            chunk_records.append({
                "text": chunk,
                "raw_id": raw_id,
                "restricts": [
                    {"namespace": "source", "allow_list": [object_name]},
                    {"namespace": "book_title", "allow_list": [book_title]},
                    {"namespace": "chapter_title", "allow_list": [chapter_title]},
                    {"namespace": "chapter_index", "allow_list": [str(chapter_index)]},
                    {"namespace": "source_file", "allow_list": [item.get_name()]},
                ],
            })

        chapter_index += 1

    logger.info(
        "epub_extracted",
        object_name=object_name,
        book_title=book_title,
        chapters=chapter_index,
        total_chunks=len(chunk_records),
    )

    return chunk_records


def _process_pdf(local_path: str, object_name: str) -> list[dict[str, Any]]:
    """Extract text from a PDF page by page and chunk each page independently.

    Args:
        local_path: Path to the downloaded PDF file.
        object_name: GCS object key (used for datapoint IDs and metadata).

    Returns:
        List of chunk records with text, raw_id, and restricts.
    """
    reader = PdfReader(local_path)
    splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
        encoding_name="cl100k_base",
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
    )

    chunk_records: list[dict[str, Any]] = []

    for page_num, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if not text.strip():
            continue

        chunks = splitter.split_text(text)

        for i, chunk in enumerate(chunks):
            raw_id = f"{object_name}::page{page_num}::chunk_{i}"
            chunk_records.append({
                "text": chunk,
                "raw_id": raw_id,
                "restricts": [
                    {"namespace": "source", "allow_list": [object_name]},
                    {"namespace": "page_number", "allow_list": [str(page_num + 1)]},
                ],
            })

    logger.info(
        "pdf_extracted",
        object_name=object_name,
        pages=len(reader.pages),
        total_chunks=len(chunk_records),
    )

    return chunk_records


def _process_docx(local_path: str, object_name: str) -> list[dict[str, Any]]:
    """Extract text from a DOCX file with heading-aware sectioning.

    Groups paragraphs under their nearest preceding heading and chunks
    each section independently.

    Args:
        local_path: Path to the downloaded DOCX file.
        object_name: GCS object key (used for datapoint IDs and metadata).

    Returns:
        List of chunk records with text, raw_id, and restricts.
    """
    doc = DocxDocument(local_path)
    splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
        encoding_name="cl100k_base",
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
    )

    # Group paragraphs into sections by heading
    sections: list[tuple[str, str]] = []  # (heading_title, section_text)
    current_heading = "Untitled Section"
    current_paragraphs: list[str] = []

    for para in doc.paragraphs:
        if para.style and para.style.name.startswith("Heading"):
            # Flush the previous section
            if current_paragraphs:
                sections.append((current_heading, "\n".join(current_paragraphs)))
                current_paragraphs = []
            current_heading = para.text.strip() or "Untitled Section"
        elif para.text.strip():
            current_paragraphs.append(para.text)

    # Flush final section
    if current_paragraphs:
        sections.append((current_heading, "\n".join(current_paragraphs)))

    chunk_records: list[dict[str, Any]] = []
    section_index = 0

    for heading_title, section_text in sections:
        chunks = splitter.split_text(section_text)

        for i, chunk in enumerate(chunks):
            raw_id = f"{object_name}::sec{section_index}::chunk_{i}"
            chunk_records.append({
                "text": chunk,
                "raw_id": raw_id,
                "restricts": [
                    {"namespace": "source", "allow_list": [object_name]},
                    {"namespace": "chapter_title", "allow_list": [heading_title]},
                    {"namespace": "section_index", "allow_list": [str(section_index)]},
                ],
            })

        section_index += 1

    logger.info(
        "docx_extracted",
        object_name=object_name,
        sections=len(sections),
        total_chunks=len(chunk_records),
    )

    return chunk_records


def _process_pptx(local_path: str, object_name: str) -> list[dict[str, Any]]:
    """Extract text from a PPTX file slide by slide.

    Args:
        local_path: Path to the downloaded PPTX file.
        object_name: GCS object key (used for datapoint IDs and metadata).

    Returns:
        List of chunk records with text, raw_id, and restricts.
    """
    prs = Presentation(local_path)
    splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
        encoding_name="cl100k_base",
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
    )

    chunk_records: list[dict[str, Any]] = []

    for slide_num, slide in enumerate(prs.slides):
        text_parts: list[str] = []
        slide_title = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        text_parts.append(para_text)

            # Capture the slide title from the title placeholder
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                slide_title = shape.text.strip()

        slide_text = "\n".join(text_parts)
        if not slide_text.strip():
            continue

        if not slide_title:
            slide_title = f"Slide {slide_num + 1}"

        chunks = splitter.split_text(slide_text)

        for i, chunk in enumerate(chunks):
            raw_id = f"{object_name}::slide{slide_num}::chunk_{i}"
            chunk_records.append({
                "text": chunk,
                "raw_id": raw_id,
                "restricts": [
                    {"namespace": "source", "allow_list": [object_name]},
                    {"namespace": "slide_number", "allow_list": [str(slide_num + 1)]},
                    {"namespace": "slide_title", "allow_list": [slide_title]},
                ],
            })

    logger.info(
        "pptx_extracted",
        object_name=object_name,
        slides=len(prs.slides),
        total_chunks=len(chunk_records),
    )

    return chunk_records


def _split_text(text: str) -> list[str]:
    """Split text into chunks of ~500 tokens with 50-token overlap."""
    splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
        encoding_name="cl100k_base",
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
    )
    return splitter.split_text(text)


def _embed_chunks(chunks: list[str]) -> list[list[float]]:
    """Embed text chunks using Vertex AI text-embedding-004.

    Batches requests to stay within API limits.
    """
    model = TextEmbeddingModel.from_pretrained(EMBEDDING_MODEL)
    all_embeddings: list[list[float]] = []

    for i in range(0, len(chunks), EMBEDDING_BATCH_SIZE):
        batch = chunks[i : i + EMBEDDING_BATCH_SIZE]
        inputs = [
            TextEmbeddingInput(text=chunk, task_type="RETRIEVAL_DOCUMENT")
            for chunk in batch
        ]
        results = model.get_embeddings(
            inputs,
            output_dimensionality=EMBEDDING_DIMENSIONS,
        )
        all_embeddings.extend([e.values for e in results])

    return all_embeddings


def _build_datapoints(
    chunk_records: list[dict[str, Any]],
    embeddings: list[list[float]],
) -> list[dict[str, Any]]:
    """Build Vector Search datapoints with deterministic IDs.

    Each datapoint ID is a hash of the raw_id from the chunk record,
    ensuring idempotent upserts on re-processing.
    """
    datapoints = []
    for record, embedding in zip(chunk_records, embeddings):
        datapoint_id = hashlib.sha256(record["raw_id"].encode()).hexdigest()[:40]
        datapoints.append(
            {
                "datapoint_id": datapoint_id,
                "feature_vector": embedding,
                "restricts": record["restricts"],
            }
        )
    return datapoints


def _upsert_vectors(
    index_id: str,
    datapoints: list[dict[str, Any]],
) -> None:
    """Upsert datapoints to the Vertex AI Vector Search index via streaming update."""
    index = MatchingEngineIndex(index_name=index_id)
    index.upsert_datapoints(datapoints=datapoints)
